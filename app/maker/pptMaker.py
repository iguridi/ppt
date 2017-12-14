from pptx import Presentation

def remove_spaces(text):
	text = text.replace('\n', ' ')
	text = text.replace('  ', ' ')
	text = text.replace('  ', ' ')
	return text.rstrip()

class Reading:
	#In charge of all text-related thsings.
	def __init__(self, title, addrs, body):
		self.title = title
		self.addrs = addrs
		self.body = body

		self.slides = [] #List of the text splited into chunks of a predifined max chars
		self.make_pretty()

	def make_pretty(self):
		self.body = remove_spaces(self.body)
		self.body = '	' +  self.body + '\n'

	def __repr__(self):
		return (self.title, self, addrs, self.body)

class Psalm(Reading):
	'''
	Is harder to format the psalm text so it has its own class and functions
	Arguments:
		* Same as his parent: Reading()
	'''
	def __init__(self, title, addrs, body):
		super().__init__(title, addrs,body)
		self.make_pretty()
		self.split_paragraphs()

	def make_pretty(self):
		self.body = remove_spaces(self.body)
		self.body = self.body.replace('R/. ', 'R. ')
		self.body = self.body.replace('. R. ', '. R.')
		self.body = self.body.replace('! R. ', '! R.')

	def split_paragraphs(self):
		'''
		Spliting the psalm`s paragraphs is needed for adding the "R." at the end of each.
		'''
		self.psalm_response = self.body[:self.body.find('***')]
		print(self.psalm_response)
		pars_pos =self.body.find('***') + 3
		self.body = self.body[pars_pos:]
		self.psalm_paragraphs = self.body.split(' R.')
		print(self.psalm_paragraphs)
		del self.psalm_paragraphs[-1]
		for n, _ in enumerate(self.psalm_paragraphs):
			self.psalm_paragraphs[n] = '	' + self.psalm_paragraphs[n]

	



class Maker:
	'''
	In charge of making and formating the slides.
	Arguments:
		* readings (list(string)): lectures themselves
		* base_ppt (string): name of the ppt with the layouts 
		* output_ppt (string): name of the  ppt to modify
		* slide_size (int): ideal maximum of characters in a slide
		* addrs (list(string)): readings bible addresses
		* date (string): Date of the mass
		* ppt_title (string): title of the presentation based on the celebration on that *date
	'''
	def __init__(self, readings, base_ppt, output_ppt, slide_size, addrs, date, ppt_title):
		self.readings = readings
		self.slide_size = slide_size
		self.addrs = addrs
		self.date = date
		self.ppt_title = ppt_title

		self.prs = Presentation(base_ppt)
		self.process() #All the actions
		self.prs.save(output_ppt)

	def process(self):
		'''
		The process of making and formating the ppt.
		It calls almost all of the methods of the class.
		'''
		names = ['primera_lectura',
				 'salmo',
				 'segunda_lectura',
				 'evangelio'
				 ]
		self.make_cover()
		for name in names:
			if name != 'salmo':
				self.reading = Reading(name, self.addrs[name], self.readings[name])
			else:
				self.reading = Psalm(name, self.addrs[name], self.readings[name])
			self.separate_text()
			self.make_readings_slides()
		self.add_extra_slides()

	def separate_text(self):
		'''
		Split the text into various slides depending of the text length
		'''
		t = self.reading.body
		new_text = ''
		while len(t) > self.slide_size:
			diap = t[:self.slide_size]
			i = len(diap) - diap[::-1].find('.')
			new_text +=  t[:i] + '\n\n    '
			t = t[i:]
		separated_text = new_text + t
		self.reading.slides = separated_text.split('\n\n')


	def make_cover(self):
		'''
		Add the cover and its title and date to the ppt
		'''
		cover_layout = self.prs.slide_layouts[0]
		cover = self.prs.slides.add_slide(cover_layout)
		title = cover.shapes.title #placeholder idx of the address
		date = cover.placeholders[10] #placeholder idx of the body

		title.text = self.ppt_title
		date.text = self.date


	def make_readings_slides(self):
		'''
		Add the actual slides and its text and address and ending
		'''
		ppt_readings_layout = {
							   'primera_lectura': 1,
						       'salmo': 2,
						       'segunda_lectura': 3,
						       'evangelio' : 4
						       }

		n = ppt_readings_layout[self.reading.title]
		slide_layout = self.prs.slide_layouts[n]
		for i, slide_text in enumerate(self.reading.slides):
			slide = self.prs.slides.add_slide(slide_layout)
			address = slide.placeholders[12] # Placeholder idx of the address
			body = slide.placeholders[10] # Placeholder idx of the body text

			address.text = self.format_addr(self.reading.addrs) # The address is the same for all

			#Setting up the text of the body of the psalm is different beacause its formatting 
			if self.reading.title == 'salmo':
				txt_fm =  body.text_frame
				#Add the response bold text at the beggining
				resp = txt_fm.paragraphs[0]# Te default paragraph has white space. Dont have to add a new one, you have to use it 
				resp.text = self.reading.psalm_response
				font = resp.font
				font.bold = True
				for i in self.reading.psalm_paragraphs:
					#Add each paragraph of the psalm
					par = txt_fm.add_paragraph()
					par.text = i
					run = par.add_run()
					run.text = ' R.'
					font = run.font
					font.bold = True

			#Other readings
			else:
				body.text = slide_text
				if i == len(self.reading.slides)-1:
					self.add_readings_ends(slide)


			
	def format_addr(self, addr):
		addr = addr.strip()
		addr = '(' + addr + ')'
		return addr


	def add_readings_ends(self, slide):
		endings = {'primera_lectura': ('Palabra de Dios','Te alabamos Señor'),
		   		   'segunda_lectura': ('Palabra de Dios','Te alabamos Señor'),
		   		   'evangelio':       ('Palabra del Señor', 'Gloria a tí, Señor Jesus')
		   }

		dialog_padre = slide.placeholders[13]
		character = slide.placeholders[17]
		dialog_people = slide.placeholders[16]

		dialog_padre.text = endings[self.reading.title][0]
		character.text = 'R.'
		dialog_people.text = endings[self.reading.title][1]

	def add_extra_slides(self):
		'''
		Adds the pictures slides and the announcements slides required.
		'''
		pic1 = self.prs.slides.add_slide(self.prs.slide_layouts[6])
		announcements = self.prs.slides.add_slide(self.prs.slide_layouts[5])
		pic2 = self.prs.slides.add_slide(self.prs.slide_layouts[6])







