from pptx import Presentation

COVER = "portada"
FIRST_LECTURE = "primera_lectura"
PSALM = "salmo"
SECOND_LECTURE = "segunda_lectura"
GOSPEL = "evangelio"
ANNOUNCEMENTS = "anuncios"
PICTURE = "imagen"

LAYOUTS_INDEX = {
    COVER: 0,
    FIRST_LECTURE: 1,
    PSALM: 2,
    SECOND_LECTURE: 3,
    GOSPEL: 4,
    ANNOUNCEMENTS: 5,
    PICTURE: 6,
}


class Ppt:
    def __init__(self, base_ppt: str) -> str:
        self.prs = Presentation(base_ppt)

    def get_slide_layout(self, name: str) -> str:
        return self.prs.slide_layouts[LAYOUTS_INDEX[name]]

    def add_slide(self, name: str) -> str:
        return self.prs.slides.add_slide(self.get_slide_layout(name))

    def save(self, output_ppt: str) -> None:
        self.prs.save(output_ppt)
