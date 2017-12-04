from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Presentation

class Lectura:
	def __init__(self, title, text):
		pass


prs = Presentation('plantilla python.pptx')

slide = prs.slides[0]
for shape in slide.shapes:
	if not shape.has_text_frame:
		continue

	shape.text = "hoolala"



prs.save('plantilla python.pptx')