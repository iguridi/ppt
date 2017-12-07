from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

def remove_spaces(text):
	text = text.replace('\n', ' ')
	text = text.replace('  ', ' ')
	text = text.replace('  ', ' ')
	return text.rstrip()

class Reading:
	#In charge of all text-related thsings.
	def __init__(self, title, addrs, body):
		self.title = title
		self.addrs = addrs
		self.body = body

		self.slides = [] #List of the text slited into pieces of 730 chars
		self.make_pretty()

	def make_pretty(self):
		self.body = remove_spaces(self.body)
		self.body = '	' +  self.body + '\n'

class Psalm(Reading):
	def __init__(self, title, addrs, body):
		super().__init__(title, addrs,body)
		self.make_pretty()
		self.split_paragraphs()

	def make_pretty(self):
		self.body = remove_spaces(self.body)
		self.body = self.body.replace('. R. ', '. R.')
		self.body = self.body.replace('! R. ', '! R.')

	def split_paragraphs(self):
		self.psalm_paragraphs = self.body.split(' R.')
		del self.psalm_paragraphs[-1]
		for n, _ in enumerate(self.psalm_paragraphs):
			self.psalm_paragraphs[n] = '	' + self.psalm_paragraphs[n]


	


class Maker:
	
	def __init__(self, readings, base_ppt, output_ppt, slide_size, addrs, date, ppt_title):
		self.readings = readings
		self.slide_size = slide_size
		self.addrs = addrs
		self.date = date
		self.ppt_title = ppt_title

		self.prs = Presentation(base_ppt)
		self.process()
		self.prs.save(output_ppt)

	def process(self):
		names = ['primera_lectura',
				 'salmo',
				 'segunda_lectura',
				 'evangelio'
				 ]
		self.make_cover()
		for name in names:
			if name != 'salmo':
				self.reading = Reading(name, self.addrs[name], self.readings[name])
			else:
				self.reading = Psalm(name, self.addrs[name], self.readings[name])
			self.separate_text()
			self.make_readings_slides()

	def separate_text(self):
		t = self.reading.body
		new_text = ''
		while len(t) > self.slide_size:
			diap = t[:self.slide_size]
			i = len(diap) - diap[::-1].find('.')
			new_text +=  t[:i] + '\n\n    '
			t = t[i:]
		separated_text = new_text
		self.reading.slides = separated_text.split('\n\n')
		#del self.reading.slides[-1] ###FIX

	def make_cover(self):
		cover_layout = self.prs.slide_layouts[0]
		cover = self.prs.slides.add_slide(cover_layout)
		title = cover.shapes.title #placeholder idx of the address
		date = cover.placeholders[10] #placeholder idx of the body

		title.text = self.ppt_title
		date.text = self.date


	def make_readings_slides(self):
		ppt_readings_layout = {
							   'primera_lectura': 1,
						       'salmo': 2,
						       'segunda_lectura': 3,
						       'evangelio' : 4#,
						       #'oracion': 5
						       }

		n = ppt_readings_layout[self.reading.title]
		slide_layout = self.prs.slide_layouts[n]
		for i, slide_text in enumerate(self.reading.slides):
			slide = self.prs.slides.add_slide(slide_layout)
			address = slide.placeholders[12] #placeholder idx of the address
			body = slide.placeholders[10] #placeholder idx of the body

			address.text = self.reading.addrs
			body.text = slide_text

			if i == len(self.reading.slides)-1 and self.reading.title != 'salmo':
				self.add_end(self.reading.title, slide)
				endings = {'primera_lectura': ('Palabra de Dios','Te alabamos Señor'),
				   		   'segunda_lectura': ('Palabra de Dios','Te alabamos Señor'),
				   		   'evangelio':       ('Palabra del Señor', 'Gloria a tí, Señor Jesus')
				   }

				char1 = slide.placeholders[15] 
				dialog_padre = slide.placeholders[13]
				char1.text = 'L.'
				dialog_padre.text = endings[self.reading.title][0]

				char2 = slide.placeholders[17]
				dialog_people = slide.placeholders[16]
				char2.text = 'R.'
				dialog_people.text = endings[self.reading.title][1]

			elif i == len(self.reading.slides)-1 and self.reading.title == 'salmo':
				txtfm =  body.text_frame
				for paragraph in txtfm.paragraphs:

					for i in self.reading.psalm_paragraphs:
						par = txtfm.add_paragraph()
						par.text = i
						run = par.add_run()
						run.text = ' R.'
						font = run.font
						font.bold = True

				

	def add_end(self, title, slide):
		#15 = R, 13=Dialog
		pass
		
		











































