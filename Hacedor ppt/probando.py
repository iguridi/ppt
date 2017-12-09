from pptx import Presentation

prs = Presentation('plantilla python.pptx')

slide = prs.slides.add_slide(prs.slide_layouts[1])
address = slide.placeholders[12] #placeholder idx of the address
body = slide.placeholders[10] #placeholder idx of the body
address.text = '(Jn 15-45.56)'
body.text = "Lorem ipsum dolor sit amet, consectetur adipiscing elit, sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris nisi ut aliquip ex ea commodo consequat. Duis aute irure dolor in reprehenderit in voluptate velit esse cillum dolore eu fugiat nulla pariatur. Excepteur sint occaecat cupidatat non proident, sunt in culpa qui officia deserunt mollit anim id est laborum"

prs.save('probando.pptx')
